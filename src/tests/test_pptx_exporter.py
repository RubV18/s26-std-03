from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation

from futurdata_pptx import ExportOptions, PPTXExportEngine


REPOSITORY_ROOT = Path(__file__).resolve().parents[2]
SAMPLE = (
    REPOSITORY_ROOT
    / "docs"
    / "pptx-exporter"
    / "source"
    / "example_air_fryer.json"
)


def slide_text(slide) -> str:
    values = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            values.append(shape.text)
        if shape.has_table:
            for row in shape.table.rows:
                values.extend(cell.text for cell in row.cells)
    return "\n".join(values)


def test_complete_export(tmp_path: Path):
    output = tmp_path / "air_fryer.pptx"
    result, issues = PPTXExportEngine().load_validate_export(
        SAMPLE,
        output,
        ExportOptions(),
        prefer_external_loader=False,
    )
    assert result.exists()
    prs = Presentation(result)
    all_text = "\n".join(slide_text(slide) for slide in prs.slides)
    assert "bill of materials" in all_text.lower()
    assert "Plastic cover (8 hooks)" in all_text
    assert "remove the cover" in all_text
    assert "Branch" in all_text or "branch" in all_text
    assert any(issue.rule == "mass_balance" for issue in issues)
    assert any(issue.rule == "branch_source" for issue in issues)


def test_branch_source_is_not_invented(tmp_path: Path):
    engine = PPTXExportEngine()
    document = engine.load(SAMPLE, prefer_external_loader=False)
    source, method = document.source_for_step(8)  # Step 9, after Step 8 ended a branch.
    assert source is None
    assert method == "unspecified_branch"

    output = engine.export(
        document,
        tmp_path / "branch.pptx",
        ExportOptions(start_step=9, end_step=9, include_bom=False),
    )
    prs = Presentation(output)
    text = "\n".join(slide_text(slide) for slide in prs.slides)
    assert "Source: not specified in JSON" in text


def test_grouping_and_limit(tmp_path: Path):
    engine = PPTXExportEngine()
    document = engine.load(SAMPLE, prefer_external_loader=False)
    engine.validate(document)
    output = engine.export(
        document,
        tmp_path / "partial.pptx",
        ExportOptions(
            start_step=3,
            end_step=12,
            max_action_groups=6,
            groups_per_slide=2,
            include_bom=False,
            include_tools_summary=False,
            include_safety_summary=False,
        ),
    )
    prs = Presentation(output)
    # title + overview + warnings + 3 grouped slides + closing
    assert len(prs.slides) >= 7
    text = "\n".join(slide_text(slide) for slide in prs.slides)
    assert "Action groups 3–4" in text
    assert "Action groups 7–8" in text
    assert "Action groups 9–10" not in text


def test_explicit_source_and_step_mass_balance(tmp_path: Path):
    payload = json.loads(SAMPLE.read_text(encoding="utf-8"))
    payload["steps"][8]["source"] = {
        "node_id": 999,
        "name": "Top assembly",
        "weight": 250.0,
        "weight_unit": "g",
    }
    source_file = tmp_path / "explicit.json"
    source_file.write_text(json.dumps(payload), encoding="utf-8")
    engine = PPTXExportEngine()
    document = engine.load(source_file, prefer_external_loader=False)
    source, method = document.source_for_step(8)
    assert source is not None and source.name == "Top assembly"
    assert method == "explicit"
    issues = engine.validate(document)
    assert not any(issue.rule == "branch_source" and issue.location == "steps[8].source" for issue in issues)


def test_settings_roundtrip(tmp_path: Path):
    options = ExportOptions(start_step=2, end_step=8, groups_per_slide=2, include_images=False)
    path = options.save(tmp_path / "settings.json")
    loaded = ExportOptions.load(path)
    assert loaded.start_step == 2
    assert loaded.end_step == 8
    assert loaded.groups_per_slide == 2
    assert loaded.include_images is False
